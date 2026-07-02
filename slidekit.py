#!/usr/bin/env python3
"""
slidekit.py — Construtor de slides UNIFOR que replica o visual ENRIQUECIDO do
deck do Encontro 1 (não os placeholders do gerar_slides.py da skill).

API:
    from slidekit import build_deck
    build_deck(spec, out_path, assets_dir=None)

`spec` é um dict (ou caminho .json) com:
{
  "disciplina": "Linguagens de Programação para Engenharia de Dados",
  "footer": "Linguagens de Programação para Engenharia de Dados · Encontro 2 · 26/06/2026",
  "slides": [ {type, ...}, ... ]
}

Tipos de slide (espelham o E1):
  - capa:        {disciplina, curso, encontro_label, datas, carga, professor, linkedin}
  - abertura:    {modulo, encontro_titulo, horario}      # módulo + encontro centralizado, logo positiva no canto
  - conteudo:    {titulo, subtitulo?, bullets:[...]}      # barra lateral, título azul, subtítulo cinza, bullets
  - citacao:     {linhas:[...]}                           # fundo azul cheio, citação branca centralizada
  - pratica:     {bullets:[...]}                          # banner lateral grosso + "Prática"
  - bibliografia:{itens:[...]}
  - contato:     {nome, linkedin}                         # fundo azul, logo grande, "Obrigado!"

Cada slide pode trazer "footer" para sobrescrever o footer global.
Requer python-pptx.
"""
from __future__ import annotations
import json, glob, os
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- tema (E1)
TEMA = {
    "cores": {
        "primaria": "#003B7A", "primaria_escura": "#002855", "primaria_clara": "#1E5DAA",
        "secundaria": "#E30613", "fundo": "#FFFFFF", "fundo_alt": "#F4F6FA",
        "texto_principal": "#1A1A1A", "texto_secundario": "#555555",
        "destaque": "#FFC107", "divisor": "#D0D5DD",
    },
    "tipografia": {
        "titulo_familia": "Calibri", "corpo_familia": "Calibri",
        "tamanho_titulo_capa": 40, "tamanho_titulo": 32, "tamanho_subtitulo": 22,
        "tamanho_corpo": 18, "tamanho_rodape": 11,
    },
    "layout": {"largura_slide_in": 13.333, "altura_slide_in": 7.5,
               "barra_lateral_in": 0.35, "rodape_altura_in": 0.4},
    "logo": {"aspect_ratio": 2.36},
}
W = TEMA["layout"]["largura_slide_in"]
H = TEMA["layout"]["altura_slide_in"]


def _rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _find_assets() -> Path | None:
    # localiza a pasta assets da skill aula-unifor (logos), tolerante a sessão
    cands = glob.glob("/sessions/*/mnt/.claude/skills/aula-unifor/assets")
    cands += glob.glob(os.path.expanduser("~") + "/**/aula-unifor/assets", recursive=True)
    for c in cands:
        if Path(c, "logo").is_dir():
            return Path(c)
    return None


# ---------------------------------------------------------------- primitivas
def _bg_branco(slide):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = _rgb(TEMA["cores"]["fundo"])


def _bg_azul(slide):
    bloco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    bloco.fill.solid(); bloco.fill.fore_color.rgb = _rgb(TEMA["cores"]["primaria"])
    bloco.line.fill.background()


def _barra_lateral(slide):
    f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                               Inches(TEMA["layout"]["barra_lateral_in"]), Inches(H))
    f.fill.solid(); f.fill.fore_color.rgb = _rgb(TEMA["cores"]["primaria"])
    f.line.fill.background()


def _rodape(slide, texto: str):
    if not texto:
        return
    ah = Inches(TEMA["layout"]["rodape_altura_in"])
    linha = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5),
                                   Inches(H) - ah, Inches(W - 1.0), Inches(0.02))
    linha.fill.solid(); linha.fill.fore_color.rgb = _rgb(TEMA["cores"]["divisor"])
    linha.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(H) - ah + Inches(0.05),
                                  Inches(W - 1.0), Inches(0.3))
    p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = texto
    r.font.name = TEMA["tipografia"]["corpo_familia"]
    r.font.size = Pt(TEMA["tipografia"]["tamanho_rodape"])
    r.font.color.rgb = _rgb(TEMA["cores"]["texto_secundario"])


def _titulo(slide, texto: str, y=0.5):
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(W - 1.6), Inches(1.0))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = texto
    r.font.name = TEMA["tipografia"]["titulo_familia"]
    r.font.size = Pt(TEMA["tipografia"]["tamanho_titulo"]); r.font.bold = True
    r.font.color.rgb = _rgb(TEMA["cores"]["primaria"])


def _subtitulo(slide, texto: str, y=1.45):
    tx = slide.shapes.add_textbox(Inches(0.82), Inches(y), Inches(W - 1.6), Inches(0.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = texto
    r.font.name = TEMA["tipografia"]["corpo_familia"]
    r.font.size = Pt(TEMA["tipografia"]["tamanho_subtitulo"]); r.font.italic = True
    r.font.color.rgb = _rgb(TEMA["cores"]["texto_secundario"])


def _bullets(slide, items, y=1.9, x=1.0, size=None):
    size = size or TEMA["tipografia"]["tamanho_corpo"]
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(W - x - 0.8), Inches(4.7))
    tf = tx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(8)
        r = p.add_run(); r.text = f"•  {item}"
        r.font.name = TEMA["tipografia"]["corpo_familia"]
        r.font.size = Pt(size); r.font.color.rgb = _rgb(TEMA["cores"]["texto_principal"])


def _centro(slide, texto, y, size, cor=None, bold=False, italic=False):
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(W - 1.6), Inches(1.2))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = texto
    r.font.name = TEMA["tipografia"]["titulo_familia"]; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = _rgb(cor or TEMA["cores"]["texto_principal"])


def _logo(slide, assets: Path | None, variante: str, x, y, altura):
    if not assets:
        return
    fn = assets / "logo" / ("unifor_negative.png" if variante == "negative" else "unifor_positive.png")
    if not fn.is_file():
        return
    largura = altura * TEMA["logo"]["aspect_ratio"]
    slide.shapes.add_picture(str(fn), Inches(x), Inches(y), height=Inches(altura), width=Inches(largura))


# ---------------------------------------------------------------- slides
def _add(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg_branco(s); return s


def _capa(prs, assets, sp, footer):
    s = _add(prs)
    faixa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(2.0))
    faixa.fill.solid(); faixa.fill.fore_color.rgb = _rgb(TEMA["cores"]["primaria"])
    faixa.line.fill.background()
    _logo(s, assets, "negative", 0.7, 0.5, 1.0)
    if sp.get("curso"):
        tx = s.shapes.add_textbox(Inches(4.0), Inches(0.85), Inches(W - 4.5), Inches(0.6))
        p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = sp["curso"]
        r.font.name = TEMA["tipografia"]["corpo_familia"]; r.font.size = Pt(14)
        r.font.color.rgb = _rgb("#D9E2F0")
    _centro(s, sp.get("disciplina", "Disciplina"), 2.55,
            TEMA["tipografia"]["tamanho_titulo_capa"], TEMA["cores"]["primaria_escura"], bold=True)
    if sp.get("encontro_label"):
        _centro(s, sp["encontro_label"], 3.95, 22, TEMA["cores"]["primaria_clara"], bold=True)
    linhas = [f"Professor: {sp.get('professor','Cassio Pinheiro')}"]
    if sp.get("linkedin"):
        linhas.append(sp["linkedin"])
    _centro(s, "  •  ".join(linhas), 4.85, 16, TEMA["cores"]["texto_secundario"])
    info = []
    if sp.get("carga"):
        info.append(f"Carga horária: {sp['carga']}")
    if sp.get("datas"):
        info.append(f"Datas: {sp['datas']}")
    if info:
        _centro(s, "  |  ".join(info), 5.5, 14, TEMA["cores"]["texto_secundario"])
    _rodape(s, sp.get("footer", footer))


def _abertura(prs, assets, sp, footer):
    s = _add(prs); _barra_lateral(s)
    if sp.get("modulo"):
        _centro(s, sp["modulo"], 2.4, 18, TEMA["cores"]["primaria_clara"])
    _centro(s, sp.get("encontro_titulo", ""), 3.1, TEMA["tipografia"]["tamanho_titulo"],
            TEMA["cores"]["primaria"], bold=True)
    if sp.get("horario"):
        _centro(s, sp["horario"], 4.0, 14, TEMA["cores"]["texto_secundario"])
    alt = 0.45; larg = alt * TEMA["logo"]["aspect_ratio"]
    _logo(s, assets, "positive", W - larg - 0.5, H - alt - 0.6, alt)
    _rodape(s, sp.get("footer", footer))


def _conteudo(prs, sp, footer):
    s = _add(prs); _barra_lateral(s)
    _titulo(s, sp["titulo"])
    y = 1.9
    if sp.get("subtitulo"):
        _subtitulo(s, sp["subtitulo"]); y = 2.25
    _bullets(s, sp.get("bullets", []), y=y, size=sp.get("size"))
    _rodape(s, sp.get("footer", footer))


def _citacao(prs, sp, footer):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg_branco(s); _bg_azul(s)
    linhas = sp.get("linhas", [])
    n = len(linhas); start = 3.4 - (n - 1) * 0.55
    for i, ln in enumerate(linhas):
        bold = sp.get("bold_last", True) and i == n - 1
        _centro(s, ln, start + i * 1.1, 26 if not bold else 28,
                TEMA["cores"]["fundo"], bold=bold, italic=not bold)


def _pratica(prs, sp, footer):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg_branco(s)
    banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(1.2), Inches(H))
    banner.fill.solid(); banner.fill.fore_color.rgb = _rgb(TEMA["cores"]["primaria_clara"])
    banner.line.fill.background()
    tx = s.shapes.add_textbox(Inches(1.6), Inches(1.4), Inches(W - 2.4), Inches(1.0))
    p = tx.text_frame.paragraphs[0]; r = p.add_run()
    r.text = sp.get("titulo", "Prática")
    r.font.name = TEMA["tipografia"]["titulo_familia"]
    r.font.size = Pt(TEMA["tipografia"]["tamanho_titulo"]); r.font.bold = True
    r.font.color.rgb = _rgb(TEMA["cores"]["primaria"])
    _bullets(s, sp.get("bullets", []), y=2.8, x=1.6)
    _rodape(s, sp.get("footer", footer))


def _bibliografia(prs, sp, footer):
    s = _add(prs); _barra_lateral(s)
    _titulo(s, "Bibliografia")
    _bullets(s, sp.get("itens", []))
    _rodape(s, sp.get("footer", footer))


def _contato(prs, assets, sp, footer):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg_branco(s); _bg_azul(s)
    alt = 1.4; larg = alt * TEMA["logo"]["aspect_ratio"]
    _logo(s, assets, "negative", (W - larg) / 2, 0.8, alt)
    _centro(s, sp.get("titulo", "Obrigado!"), 3.0, 54, TEMA["cores"]["fundo"], bold=True)
    _centro(s, sp.get("nome", "Cassio Pinheiro"), 4.4, 22, "#E0EAF6")
    if sp.get("linkedin"):
        _centro(s, sp["linkedin"], 5.1, 16, "#C7D6EA")


_DISPATCH = {
    "capa": lambda prs, assets, sp, f: _capa(prs, assets, sp, f),
    "abertura": lambda prs, assets, sp, f: _abertura(prs, assets, sp, f),
    "conteudo": lambda prs, assets, sp, f: _conteudo(prs, sp, f),
    "citacao": lambda prs, assets, sp, f: _citacao(prs, sp, f),
    "pratica": lambda prs, assets, sp, f: _pratica(prs, sp, f),
    "bibliografia": lambda prs, assets, sp, f: _bibliografia(prs, sp, f),
    "contato": lambda prs, assets, sp, f: _contato(prs, assets, sp, f),
}


def build_deck(spec, out_path, assets_dir=None):
    if isinstance(spec, (str, Path)):
        spec = json.loads(Path(spec).read_text(encoding="utf-8"))
    assets = Path(assets_dir) if assets_dir else _find_assets()
    prs = Presentation()
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    footer = spec.get("footer", "")
    for sl in spec["slides"]:
        t = sl["type"]
        _DISPATCH[t](prs, assets, sl, footer)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return len(spec["slides"])


if __name__ == "__main__":
    import sys
    n = build_deck(sys.argv[1], sys.argv[2])
    print(f"[OK] {n} slides -> {sys.argv[2]}")
